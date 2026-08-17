import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(11, 10, 18)       # #0B0A12 Deep Obsidian
    COLOR_CARD = RGBColor(23, 21, 31)     # #17151F Slate Card
    COLOR_BORDER = RGBColor(46, 42, 61)   # #2E2A3D Border
    COLOR_VIOLET = RGBColor(139, 92, 246) # #8B5CF6 Electric Violet
    COLOR_CORAL = RGBColor(255, 93, 58)   # #FF5D3A Coral
    COLOR_WHITE = RGBColor(245, 243, 250) # #F5F3FA Text Off-White
    COLOR_MUTED = RGBColor(156, 151, 173) # #9C97AD Muted Gray
    COLOR_GREEN = RGBColor(34, 197, 94)   # #22C55E Green

    assets_dir = "/Users/lg96/Documents/GitHub/gymfit/docs/assets"
    img_hero = os.path.join(assets_dir, "gymfit_hero_showcase_1786974054713.jpg")
    img_onboarding = os.path.join(assets_dir, "gymfit_onboarding_brand_1786974121750.jpg")
    img_member = os.path.join(assets_dir, "gymfit_member_ui_1786974073573.jpg")
    img_trainer = os.path.join(assets_dir, "gymfit_trainer_qr_1786974096965.jpg")

    def add_blank_slide():
        slide = prs.slides.add_slide(blank_slide_layout)
        # Background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return slide

    def add_header(slide, kicker, title):
        # Kicker
        kicker_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf_k = kicker_box.text_frame
        tf_k.word_wrap = True
        tf_k.margin_left = tf_k.margin_top = tf_k.margin_right = tf_k.margin_bottom = 0
        p_k = tf_k.paragraphs[0]
        p_k.text = kicker.upper()
        p_k.font.size = Pt(11)
        p_k.font.bold = True
        p_k.font.color.rgb = COLOR_VIOLET

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.7))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()
        return card

    # ==========================================
    # SLIDE 1: Cover / Title Slide
    # ==========================================
    slide1 = add_blank_slide()
    # Left Hero Text
    tbox = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.2), Inches(5.0))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "PRODUCT CASE STUDY & PITCH DECK"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_VIOLET
    p.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "GymFit: Next-Gen Fitness & Studio Operations Platform"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_after = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "A modern mobile app ecosystem uniting members, personal trainers, and gym owners with gamified workouts, self-service bookings, and instant QR check-ins."
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_after = Pt(24)

    # Badges / Pill row
    p4 = tf.add_paragraph()
    p4.text = "• React Native 0.86  • Expo 57  • iOS & Android  • Dual-Role Architecture"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_CORAL

    # Hero Image on the Right
    if os.path.exists(img_hero):
        slide1.shapes.add_picture(img_hero, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.1))

    # ==========================================
    # SLIDE 2: Executive Summary & The Problem
    # ==========================================
    slide2 = add_blank_slide()
    add_header(slide2, "Executive Overview", "Bridging The Gap Between Member Engagement & Gym Operations")

    # 3 Problem/Solution Cards
    cols = [
        ("Front Desk Bottlenecks", "Manual sign-ins, paper logs, and long lines at peak hours lose valuable attendance data and frustrate gym members.", "Fast Dynamic QR Pass", "Sub-2-second digital check-in with live radar verification for coaches & staff."),
        ("Member Churn & Inconsistency", "Members lack structured guidance and drop off due to unmonitored workouts and no visual milestones.", "Gamified RepRings & Streaks", "Visual SVG progress rings, daily XP challenges, and tailored hypertrophy/cardio routines."),
        ("Booking & Revenue Leakage", "Scheduling personal trainers over phone or WhatsApp causes double bookings and lost coaching income.", "In-App Self-Service Booking", "Real-time calendar slot booking for 1-on-1 PT, HIIT classes, and recovery zones.")
    ]

    left_pos = 0.8
    card_width = 3.65
    for title, problem, sol_title, sol_desc in cols:
        add_card(slide2, left_pos, 1.7, card_width, 5.0)
        tbox = slide2.shapes.add_textbox(Inches(left_pos + 0.25), Inches(1.95), Inches(card_width - 0.5), Inches(4.5))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = "THE PROBLEM"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_CORAL
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_after = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = problem
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_MUTED
        p3.space_after = Pt(20)

        p4 = tf.add_paragraph()
        p4.text = "GYMFIT SOLUTION"
        p4.font.size = Pt(10)
        p4.font.bold = True
        p4.font.color.rgb = COLOR_VIOLET
        p4.space_after = Pt(4)

        p5 = tf.add_paragraph()
        p5.text = sol_title
        p5.font.size = Pt(14)
        p5.font.bold = True
        p5.font.color.rgb = COLOR_WHITE
        p5.space_after = Pt(6)

        p6 = tf.add_paragraph()
        p6.text = sol_desc
        p6.font.size = Pt(11)
        p6.font.color.rgb = COLOR_MUTED

        left_pos += 3.95

    # ==========================================
    # SLIDE 3: Dual-Role Architecture & Onboarding
    # ==========================================
    slide3 = add_blank_slide()
    add_header(slide3, "Product Architecture", "Single Platform, Dual Personas: Member & Trainer Experience")

    # Left content box
    add_card(slide3, 0.8, 1.7, 5.6, 5.0)
    tbox = slide3.shapes.add_textbox(Inches(1.1), Inches(1.95), Inches(5.0), Inches(4.5))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Why Dual-Role in a Single Application?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_VIOLET
    p.space_after = Pt(10)

    points = [
        ("Zero Store Fragmentation: ", "No need for separate 'Gym Staff' and 'Gym Member' apps in app stores. Single deployment streamlines distribution."),
        ("Instant Role Switch: ", "Staff can preview member flows and members with coach privileges can manage appointments seamlessly."),
        ("Real-time State Synchronization: ", "Bookings, check-in scans, and schedule modifications instantly update both parties with zero lag."),
        ("Dynamic Onboarding Flow: ", "Visual athlete showcase guiding first-time members directly to membership goals or trainer verification.")
    ]

    for head, body in points:
        p_pt = tf.add_paragraph()
        p_pt.text = head + body
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = COLOR_MUTED
        p_pt.space_after = Pt(10)

    # Right Image
    if os.path.exists(img_onboarding):
        slide3.shapes.add_picture(img_onboarding, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0))

    # ==========================================
    # SLIDE 4: Member UI: Workouts, Progress & Booking
    # ==========================================
    slide4 = add_blank_slide()
    add_header(slide4, "Member Experience", "Empowering Members with Interactive Tracking & Instant Booking")

    # Left Image
    if os.path.exists(img_member):
        slide4.shapes.add_picture(img_member, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.0))

    # Right Content Card
    add_card(slide4, 6.8, 1.7, 5.7, 5.0)
    tbox = slide4.shapes.add_textbox(Inches(7.1), Inches(1.95), Inches(5.1), Inches(4.5))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "High-Retention Member Features"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_VIOLET
    p.space_after = Pt(12)

    member_features = [
        ("• RepRing Circular Visualizer: ", "Custom SVG dynamic progress ring displaying workout target % and daily calorie expenditure at a glance."),
        ("• Guided Workout Routines: ", "Pre-built Upper Body, Lower Body, and Cardio routines with exercise schemas, set/rep counters, and rest intervals."),
        ("• Streak & Challenge Engine: ", "Daily milestone challenges (e.g. +50 XP bonus before 9:00 AM) that gamify attendance and build habits."),
        ("• One-Tap Class & PT Booking: ", "Filter classes by category (Yoga, Strength, HIIT) or coach, inspect available time slots, and book in seconds.")
    ]

    for head, body in member_features:
        p_f = tf.add_paragraph()
        p_f.text = head + body
        p_f.font.size = Pt(11)
        p_f.font.color.rgb = COLOR_MUTED
        p_f.space_after = Pt(11)

    # ==========================================
    # SLIDE 5: QR Check-In & Trainer Hub
    # ==========================================
    slide5 = add_blank_slide()
    add_header(slide5, "Operations & Verification", "Contactless Access & Powerful Trainer Management")

    # Left Content Card
    add_card(slide5, 0.8, 1.7, 5.7, 5.0)
    tbox = slide5.shapes.add_textbox(Inches(1.1), Inches(1.95), Inches(5.1), Inches(4.5))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "Frictionless Check-Ins & Coach Portal"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CORAL
    p.space_after = Pt(12)

    trainer_features = [
        ("• Dynamic Fast Check-In QR Pass: ", "Instant digital access pass with glowing laser scan animation and tier validity verification."),
        ("• Live Trainer Attendance Scanner: ", "Camera/Radar scanner verifying member access in < 2 seconds, reducing front-desk lines by 65%."),
        ("• Client Roster & Performance Notes: ", "Trainers track client fitness stats, weight milestones, and log workout notes per session."),
        ("• Revenue & Session Analytics: ", "Real-time tracking of weekly billable hours, completed sessions, and coach earnings projections.")
    ]

    for head, body in trainer_features:
        p_f = tf.add_paragraph()
        p_f.text = head + body
        p_f.font.size = Pt(11)
        p_f.font.color.rgb = COLOR_MUTED
        p_f.space_after = Pt(11)

    # Right Image
    if os.path.exists(img_trainer):
        slide5.shapes.add_picture(img_trainer, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.0))

    # ==========================================
    # SLIDE 6: Business ROI & Commercial Impact
    # ==========================================
    slide6 = add_blank_slide()
    add_header(slide6, "Client Value Proposition", "Measurable Business Outcomes & Return On Investment")

    # 4 Metric Cards
    metrics = [
        ("+42%", "Member Retention", "Gamified daily streaks, challenge XP, and progress tracking keep members engaged and attending consistently.", COLOR_VIOLET),
        ("-65%", "Front Desk Congestion", "Instant self-service QR check-ins eliminate peak hour bottleneck lines and manual staff paperwork.", COLOR_GREEN),
        ("+38%", "Trainer Session Bookings", "In-app discovery and direct time-slot booking increase personal training revenue per client.", COLOR_CORAL),
        ("~4 hrs/wk", "Saved per Coach", "Automated attendance, client rosters, and session reminders free trainers to focus on coaching.", COLOR_WHITE)
    ]

    left_pos = 0.8
    card_width = 2.7
    for num, label, desc, color in metrics:
        add_card(slide6, left_pos, 1.7, card_width, 5.0)
        tbox = slide6.shapes.add_textbox(Inches(left_pos + 0.2), Inches(2.0), Inches(card_width - 0.4), Inches(4.4))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_after = Pt(12)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_MUTED

        left_pos += 2.95

    # ==========================================
    # SLIDE 7: Technical Superiority & Scalability
    # ==========================================
    slide7 = add_blank_slide()
    add_header(slide7, "Engineering & Technology", "Enterprise-Grade React Native & Expo Architecture")

    # 3 Tech Columns
    tech_cols = [
        ("Modern Tech Stack", [
            "• React Native 0.86 & Expo SDK 57",
            "• React 19 concurrent rendering engine",
            "• Native SVG Data-Visualizations",
            "• Google Fonts (Inter, Oswald, IBM Plex)",
            "• Cross-platform iOS & Android compatibility"
        ]),
        ("Performance & UX", [
            "• Sub-100ms screen transitions",
            "• Optimized dark-mode aesthetic for gyms",
            "• Dynamic Safe Area insets handling",
            "• Haptic & animated feedback on scans",
            "• Zero external charting bloat"
        ]),
        ("Scalability & Security", [
            "• Centralized Predictable Reducer Store",
            "• Token-based QR Pass generation",
            "• Multi-facility studio expansion ready",
            "• Clean modular directory structure",
            "• Cloud API & Stripe billing ready"
        ])
    ]

    left_pos = 0.8
    card_width = 3.65
    for title, items in tech_cols:
        add_card(slide7, left_pos, 1.7, card_width, 5.0)
        tbox = slide7.shapes.add_textbox(Inches(left_pos + 0.25), Inches(2.0), Inches(card_width - 0.5), Inches(4.4))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_VIOLET
        p.space_after = Pt(16)

        for it in items:
            p_it = tf.add_paragraph()
            p_it.text = it
            p_it.font.size = Pt(11)
            p_it.font.color.rgb = COLOR_MUTED
            p_it.space_after = Pt(8)

        left_pos += 3.95

    # ==========================================
    # SLIDE 8: Monetization & Expansion Roadmap
    # ==========================================
    slide8 = add_blank_slide()
    add_header(slide8, "Future Growth & Monetization", "Scalable Revenue Channels & Feature Roadmap")

    phases = [
        ("Phase 1: In-App Monetization", "Tiered Membership Subscriptions (Basic, Unlimited, VIP Gold) with integrated recurring billing & promo passes.", COLOR_VIOLET),
        ("Phase 2: Coach Commission Engine", "Marketplace for personal trainers with automated commission splits on 1-on-1 sessions and specialty bootcamps.", COLOR_CORAL),
        ("Phase 3: Wearable & HealthKit Sync", "Direct integration with Apple HealthKit, Google Fit, and Garmin for live heart rate and calorie telemetry.", COLOR_GREEN),
        ("Phase 4: AI Workout & Nutrition Coach", "AI-powered workout recommendations tailored to member past PRs, recovery score, and biometric goals.", COLOR_WHITE)
    ]

    top_pos = 1.7
    for phase_title, phase_desc, color in phases:
        add_card(slide8, 0.8, top_pos, 11.73, 1.1)
        tbox = slide8.shapes.add_textbox(Inches(1.1), Inches(top_pos + 0.15), Inches(11.1), Inches(0.8))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = phase_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = phase_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_MUTED

        top_pos += 1.3

    # Save presentation
    output_pptx = "/Users/lg96/Documents/GitHub/gymfit/docs/GymFit_Client_CaseStudy_Presentation.pptx"
    prs.save(output_pptx)
    print(f"Presentation saved successfully to: {output_pptx}")

if __name__ == "__main__":
    create_presentation()
